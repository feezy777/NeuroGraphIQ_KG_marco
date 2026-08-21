"""
Generate NeuroGraphIQ KG V3 knowledge graph construction presentation.
Requires: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x36, 0x5D)
MID_BLUE    = RGBColor(0x31, 0x82, 0xCE)
LIGHT_BLUE  = RGBColor(0xEB, 0xF4, 0xFF)
ORANGE      = RGBColor(0xDD, 0x6B, 0x20)
GREEN       = RGBColor(0x38, 0xA1, 0x69)
AMBER       = RGBColor(0xD6, 0x9E, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x2D, 0x37, 0x42)
MID_GRAY    = RGBColor(0x71, 0x80, 0x96)
LIGHT_GRAY  = RGBColor(0xF5, 0xF6, 0xF8)
RED_ACCENT  = RGBColor(0xE5, 0x3E, 0x3E)

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helper Functions ───────────────────────────────────────────

def add_bg(slide, color=LIGHT_GRAY):
    """Solid background fill."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=Pt(14),
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    """Add a text box with single-style text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.auto_size = None
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    try:
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    except Exception:
        pass
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=Pt(13),
                  color=DARK_GRAY, bold_first=False, spacing=Pt(4),
                  font_name='Microsoft YaHei'):
    """Add a text box with multiple paragraphs. lines is a list of (text, optional_is_bold)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld = item, False
        else:
            txt, bld = item[0], item[1] if len(item) > 1 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bld or (bold_first and i == 0)
        p.font.name = font_name
        p.space_after = spacing
    return txBox

def slide_number(slide, num):
    """Add a small slide number at bottom-right."""
    add_textbox(slide, W - Inches(1.2), H - Inches(0.55), Inches(1), Inches(0.4),
                str(num), font_size=Pt(10), color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def title_bar(slide, title, subtitle=None):
    """Standard page header: dark blue bar + title + optional subtitle."""
    bar = add_rect(slide, 0, 0, W, Inches(1.15), fill_color=DARK_BLUE)
    add_textbox(slide, Inches(0.7), Inches(0.18), W - Inches(1.4), Inches(0.55),
                title, font_size=Pt(28), color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(0.68), W - Inches(1.4), Inches(0.35),
                    subtitle, font_size=Pt(13), color=RGBColor(0xA0, 0xC4, 0xE8))

def accent_line(slide, left, top, width):
    """Thin orange accent line."""
    add_rect(slide, left, top, width, Pt(3), fill_color=ORANGE)

def add_table(slide, left, top, col_widths, headers, rows, font_size=Pt(11)):
    """Add a styled table. headers: list of str. rows: list of list of str."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, tbl_width,
                                         Inches(0.35) * n_rows)
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw
    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.color.rgb = DARK_GRAY
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape

def process_card(slide, left, top, width, height, title, body_lines, color=MID_BLUE):
    """Styled content card with colored left border."""
    # Left accent bar
    add_rect(slide, left, top, Pt(5), height, fill_color=color)
    # Card background
    add_rect(slide, left + Pt(5), top, width - Pt(5), height, fill_color=WHITE)
    # Title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.3),
                title, font_size=Pt(15), color=color, bold=True)
    # Body
    add_multiline(slide, left + Inches(0.2), top + Inches(0.4), width - Inches(0.4),
                  height - Inches(0.5), body_lines, font_size=Pt(11), color=DARK_GRAY)

def flow_arrow(slide, left, top, width, color=MID_BLUE):
    """Right-pointing arrow for flow diagrams."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Pt(18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def flow_box(slide, left, top, width, height, text, color=MID_BLUE, font_size=Pt(10)):
    """Rounded box for flow diagrams."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl, DARK_BLUE)

# Large centered title
add_textbox(sl, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
            'NeuroGraphIQ KG V3', font_size=Pt(44), color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1.5), Inches(2.9), Inches(10.3), Inches(0.7),
            '多粒度脑区知识图谱构建', font_size=Pt(30), color=RGBColor(0xA0, 0xC4, 0xE8),
            bold=False, alignment=PP_ALIGN.CENTER)

accent_line(sl, Inches(5), Inches(3.65), Inches(3.3))

add_textbox(sl, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.6),
            '从脑图谱资源到可探索知识图谱的全流程自动化系统',
            font_size=Pt(16), color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.5),
            '2026 年 8 月', font_size=Pt(14), color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom decorative line
add_rect(sl, Inches(1.5), Inches(7.0), Inches(10.3), Pt(1), fill_color=RGBColor(0x3A, 0x56, 0x7D))


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Problem, Motivation & Scale
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, '问题、动机与规模')
slide_number(sl, 2)

# Problem boxes — top half
problems = [
    ('脑图谱资源分散', 'AAL3、Brainnetome、HCP-MMP、Allen\n等 6+ 种图谱，格式各异 (XML/xlsx/OWL)'),
    ('缺乏结构化知识连接', '不同粒度脑区之间无统一语义关联，\n跨图谱查询困难'),
    ('LLM 提取需治理', '大模型能做知识提取，但缺乏治理\n框架会传播错误到正式库'),
]
for i, (title, desc) in enumerate(problems):
    left = Inches(0.5 + i * 4.2)
    process_card(sl, left, Inches(1.5), Inches(3.9), Inches(1.85), title,
                 [desc], color=RED_ACCENT if i == 2 else ORANGE)

# Scale — bottom half
add_textbox(sl, Inches(0.7), Inches(3.7), Inches(4), Inches(0.35),
            '项目规模', font_size=Pt(16), color=DARK_BLUE, bold=True)

scale_data = [
    ['42 路由 · 88 服务', 'React 18 + Vite', 'PostgreSQL', '1,173 函数'],
    ['FastAPI · Python 3.11+', '14 个核心页面', '5 Schema · 59 迁移', '76 测试文件'],
]
top = Inches(4.1)
for ri, row in enumerate(scale_data):
    for ci, val in enumerate(row):
        left = Inches(0.5 + ci * 3.2)
        box = add_rect(sl, left, top + Inches(ri * 1.3), Inches(3.0), Inches(1.1), fill_color=WHITE)
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        box.line.width = Pt(0.5)
        icon_map = ['📡', '🖥️', '🗄️', '✅']
        add_textbox(sl, left + Inches(0.2), top + Inches(ri * 1.3) + Inches(0.1),
                    Inches(2.6), Inches(0.35),
                    icon_map[ci] + '  ' + val.split('·')[0].strip() if '·' in val else val,
                    font_size=Pt(12), color=DARK_BLUE, bold=True)
        if '·' in val or '\n' in val:
            parts = val.split('·')
            sub = parts[1].strip() if len(parts) > 1 else ''
        else:
            sub = ''
        if not sub and '\n' not in val:
            sub = row[ci] if ri == 0 else ''
        if sub:
            add_textbox(sl, left + Inches(0.2), top + Inches(ri * 1.3) + Inches(0.55),
                        Inches(2.6), Inches(0.4), sub,
                        font_size=Pt(11), color=MID_GRAY)

# Bottom: solution
add_textbox(sl, Inches(0.7), Inches(6.8), Inches(11.5), Inches(0.4),
            '▶ 解决方案：分层漏斗治理 —— 候选 → 校验 → LLM提取 → Mirror KG → Human Review → Final KG，兼顾自动化与人工质控',
            font_size=Pt(12), color=DARK_BLUE, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Knowledge Layers & Build Pipeline
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, '知识体系与构建流水线')
slide_number(sl, 3)

# 7 layers as horizontal flow
layers = ['脑区实体', '连接', '回路', '功能', '证据', '三元组', '映射']
layer_colors = [DARK_BLUE, MID_BLUE, RGBColor(0x2B, 0x6C, 0xB0), ORANGE, AMBER, GREEN, MID_GRAY]
lx = Inches(0.4)
for i, (name, lc) in enumerate(zip(layers, layer_colors)):
    flow_box(sl, lx, Inches(1.4), Inches(1.65), Inches(0.55), name, color=lc, font_size=Pt(11))
    if i < len(layers) - 1:
        flow_arrow(sl, lx + Inches(1.7), Inches(1.52), Inches(0.15), color=lc)
    lx += Inches(1.85)

# Pipeline
add_textbox(sl, Inches(0.7), Inches(2.3), Inches(4), Inches(0.3),
            '构建流水线', font_size=Pt(16), color=DARK_BLUE, bold=True)

# Flow boxes row 1
flow_items = [
    ('资源登记', DARK_BLUE), ('批次导入', DARK_BLUE), ('原始解析', MID_BLUE), ('候选生成', MID_BLUE),
    ('规则校验', ORANGE),
]
fx = Inches(0.4)
for name, fc in flow_items:
    flow_box(sl, fx, Inches(2.8), Inches(2.2), Inches(0.55), name, color=fc, font_size=Pt(11))
    flow_arrow(sl, fx + Inches(2.3), Inches(2.92), Inches(0.15), color=fc)
    fx += Inches(2.45)

# Flow boxes row 2
flow_items2 = [
    ('LLM提取\n(DeepSeek+Kimi)', RGBColor(0x80, 0x50, 0xB0)),
    ('Mirror KG\n(预正式层)', RGBColor(0x6B, 0x46, 0xA0)),
    ('校验中心\n(三道闸门)', RED_ACCENT),
    ('人工审核', RGBColor(0xC0, 0x56, 0x20)),
    ('晋升→Final KG', GREEN),
]
fx = Inches(0.4)
for name, fc in flow_items2:
    flow_box(sl, fx, Inches(3.65), Inches(2.2), Inches(0.7), name, color=fc, font_size=Pt(10))
    if name != '晋升→Final KG':
        flow_arrow(sl, fx + Inches(2.3), Inches(3.85), Inches(0.15), color=fc)
    fx += Inches(2.45)

# Write boundary table
add_textbox(sl, Inches(0.7), Inches(4.7), Inches(4), Inches(0.3),
            '核心写边界', font_size=Pt(14), color=DARK_BLUE, bold=True)
add_table(sl, Inches(0.7), Inches(5.1),
          [Inches(2.2), Inches(3.8), Inches(5.5)],
          ['阶段', '可写入', '严禁写入'],
          [
              ['LLM 提取', 'mirror_*, llm_extraction_*', 'final_*, 自动审核, 自动晋升'],
              ['人工审核', '审核记录, 编辑建议', 'final_* 直接写入'],
              ['Promotion', 'final_* + 审计日志', '绕过审核环节'],
          ], font_size=Pt(11))

# Key principle box
add_rect(sl, Inches(0.7), Inches(6.5), Inches(11.5), Inches(0.5), fill_color=DARK_BLUE)
add_textbox(sl, Inches(1.0), Inches(6.52), Inches(10.9), Inches(0.45),
            '核心原则：每阶段向上游隔离 · LLM 是提取工具而非终审官 · 人工审核是进入 Final KG 的唯一闸门',
            font_size=Pt(13), color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Data Import & Candidate Generation
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, '数据导入与候选生成')
slide_number(sl, 4)

# Dual pipeline boxes
add_textbox(sl, Inches(0.7), Inches(1.5), Inches(5), Inches(0.3),
            '双链路独立导入', font_size=Pt(16), color=DARK_BLUE, bold=True)

add_table(sl, Inches(0.7), Inches(1.95),
          [Inches(1.4), Inches(1.8), Inches(1.0), Inches(1.2), Inches(2.0), Inches(2.6)],
          ['图谱', '粒度层', '格式', '规模', '解析器', '候选生成器'],
          [
              ['AAL3', 'macro_clinical', 'XML', '166 ROI', 'aal3_xml', 'generate-candidates'],
              ['Macro96', 'macro_clinical', 'Excel', '96 脑区', 'macro96_xlsx', 'generate-macro96-candidates'],
          ], font_size=Pt(12))

# Key mechanisms
add_textbox(sl, Inches(0.7), Inches(3.2), Inches(5), Inches(0.3),
            '关键技术机制', font_size=Pt(16), color=DARK_BLUE, bold=True)

mechanisms = [
    ('导入批次 (Import Batch)', '核心追踪单元，记录文件绑定、解析器兼容性自动检查、事件日志 (queue→start→complete/cancel)'),
    ('幂等解析', '唯一索引防止重复解析，支持安全重跑'),
    ('候选溯源', '每条候选记录 source_atlas · source_version · import_batch_id · resource_id'),
    ('候选池 (Candidate Pool)', '跨批次汇总候选数据，支持按粒度/图谱/批次筛选，为 LLM 批量提取提供数据基础'),
]
y = Inches(3.6)
for title, desc in mechanisms:
    add_textbox(sl, Inches(0.9), y, Inches(2.5), Inches(0.3),
                '▸ ' + title, font_size=Pt(12), color=DARK_BLUE, bold=True)
    add_textbox(sl, Inches(3.6), y, Inches(9.0), Inches(0.55),
                desc, font_size=Pt(11), color=DARK_GRAY)
    y += Inches(0.65)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Rule Validation & Enhancement Engine
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, '规则校验与数据增强引擎')
slide_number(sl, 5)

# Validation rules
add_textbox(sl, Inches(0.7), Inches(1.5), Inches(5), Inches(0.3),
            '12 条确定性校验规则（不依赖 LLM）', font_size=Pt(16), color=DARK_BLUE, bold=True)

add_table(sl, Inches(0.7), Inches(1.95),
          [Inches(2.5), Inches(5.0), Inches(1.5)],
          ['检查类别', '检查项', '级别'],
          [
              ['完整性', '必填字段非空 (name, source_atlas, granularity_level)', 'BLOCKER'],
              ['语义ID', 'semantic_id 格式合法性', 'BLOCKER'],
              ['唯一性', '同图谱内候选脑区不重复', 'BLOCKER'],
              ['拓扑', '脑区间引用有效性', 'BLOCKER'],
              ['溯源', 'source_atlas / version / resource_id 齐全', 'WARNING'],
              ['证据', 'evidence_text 完整度', 'WARNING'],
          ], font_size=Pt(11))

# Quality Score
add_textbox(sl, Inches(0.7), Inches(4.5), Inches(5), Inches(0.3),
            'Quality Score (0-100 加权评分)', font_size=Pt(14), color=DARK_BLUE, bold=True)
add_textbox(sl, Inches(0.7), Inches(4.85), Inches(11.5), Inches(0.3),
            '完整性 30%  +  溯源 20%  +  拓扑 20%  +  证据 20%  +  区域关联 10%',
            font_size=Pt(12), color=MID_GRAY)

# Enhancement Engine
add_textbox(sl, Inches(0.7), Inches(5.3), Inches(5), Inches(0.3),
            '数据增强引擎', font_size=Pt(16), color=DARK_BLUE, bold=True)

# Tier 1 card
process_card(sl, Inches(0.7), Inches(5.7), Inches(5.7), Inches(1.3),
             'Tier 1: 确定性自动修复',
             ['补充缺失字段 · 标准化名称 · 修复引用', '不调 LLM — 零成本、确定性、即时反馈'],
             color=GREEN)
# Tier 2 card
process_card(sl, Inches(6.8), Inches(5.7), Inches(5.7), Inches(1.3),
             'Tier 2: LLM 辅助增强 (DeepSeek)',
             ['分析疑难问题 → 生成修复建议 → mirror_enhancement_suggestions', '人工 approve / reject — 始终保留人工决策权'],
             color=MID_BLUE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — LLM Extraction Capabilities
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, 'LLM 提取能力全景')
slide_number(sl, 6)

# Dual LLM architecture
add_textbox(sl, Inches(0.7), Inches(1.5), Inches(5), Inches(0.3),
            '双 LLM 架构', font_size=Pt(16), color=DARK_BLUE, bold=True)
add_textbox(sl, Inches(0.7), Inches(1.85), Inches(11.5), Inches(0.5),
            'DeepSeek (v4-pro / V3 / R1)  +  Kimi (Moonshot)  ·  Provider 抽象层  ·  API Key 前端不可见  ·  完整 audit trail',
            font_size=Pt(12), color=MID_GRAY)

# 7 capabilities — tree layout
# Center: Candidate Regions
flow_box(sl, Inches(5.2), Inches(2.5), Inches(2.8), Inches(0.6),
         '候选脑区实体 (Regions)', color=DARK_BLUE, font_size=Pt(13))

# Level 1: 3 branches
l1_items = [('连接提取', Inches(0.5)), ('功能提取', Inches(5.2)), ('回路提取', Inches(9.8))]
for name, lx in l1_items:
    flow_box(sl, lx, Inches(3.5), Inches(2.8), Inches(0.55), name, color=MID_BLUE, font_size=Pt(12))

# Vertical connectors from center
for lx in [Inches(1.9), Inches(6.6), Inches(11.2)]:
    add_rect(sl, lx, Inches(3.1), Pt(3), Inches(0.45), fill_color=MID_BLUE)

# Level 2: derived extractions
l2_data = [
    ('投射功能提取\n(Projection Functions)', Inches(0.5), Inches(4.5)),
    ('回路功能提取\n(Circuit Functions)', Inches(5.1), Inches(4.5)),
    ('回路步骤提取\n(Circuit Steps)', Inches(9.7), Inches(4.5)),
]
for name, lx, ly in l2_data:
    flow_box(sl, lx, ly, Inches(2.9), Inches(0.75), name, color=ORANGE, font_size=Pt(11))
    # vertical connector
    add_rect(sl, lx + Inches(1.4), ly - Inches(0.4), Pt(3), Inches(0.42), fill_color=ORANGE)

# Level 3: Triple consolidation
flow_box(sl, Inches(5.2), Inches(5.7), Inches(2.8), Inches(0.55),
         '三元组整合 (Triples)', color=GREEN, font_size=Pt(12))
add_rect(sl, Inches(6.6), Inches(5.25), Pt(3), Inches(0.5), fill_color=GREEN)

# Key note
add_textbox(sl, Inches(0.7), Inches(6.6), Inches(11.5), Inches(0.4),
            '⚠ 全部在同粒度内操作 · 跨粒度关系需显式 Mapping 表 · 三元组整合为确定性转换（不调 LLM）',
            font_size=Pt(12), color=ORANGE, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Composite Workflow & Mirror Write Path
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, '复合工作流与 Mirror 写入链路')
slide_number(sl, 7)

# Composite Workflow
add_textbox(sl, Inches(0.7), Inches(1.5), Inches(5), Inches(0.3),
            '复合工作流 (Composite Workflow)', font_size=Pt(16), color=DARK_BLUE, bold=True)

steps = [
    ('Step 1', '连接 + 功能提取', 'pack 1 … N', MID_BLUE),
    ('Step 2', '回路 + 步骤提取', 'pack 1 … N', RGBColor(0x2B, 0x6C, 0xB0)),
    ('Step 3', '投射提取', 'pack 1 … N', ORANGE),
    ('Step 4', '三元组整合', '确定性，不调 LLM', GREEN),
]
sx = Inches(0.7)
for label, desc, sub, sc in steps:
    add_rect(sl, sx, Inches(2.0), Inches(2.8), Inches(1.6), fill_color=WHITE)
    add_textbox(sl, sx + Inches(0.15), Inches(2.05), Inches(2.5), Inches(0.25),
                label, font_size=Pt(11), color=sc, bold=True)
    add_textbox(sl, sx + Inches(0.15), Inches(2.3), Inches(2.5), Inches(0.35),
                desc, font_size=Pt(12), color=DARK_BLUE, bold=True)
    add_textbox(sl, sx + Inches(0.15), Inches(2.65), Inches(2.5), Inches(0.25),
                sub, font_size=Pt(10), color=MID_GRAY)
    # Bottom accent
    add_rect(sl, sx, Inches(3.55), Inches(2.8), Pt(4), fill_color=sc)
    if label != 'Step 4':
        flow_arrow(sl, sx + Inches(2.85), Inches(2.7), Inches(0.2), color=sc)
    sx += Inches(3.1)

# Features
features_text = '· Pack 机制 (pairs_per_pack 可调, 默认 20)    · Dry Run 预览 (pack 数 / token 量 / 费用估算)    · Skip Existing (避免重复)    · 暂停 / 取消 / 恢复'
add_textbox(sl, Inches(0.7), Inches(3.85), Inches(12), Inches(0.3),
            features_text, font_size=Pt(11), color=MID_GRAY)

# Mirror write path
add_textbox(sl, Inches(0.7), Inches(4.3), Inches(5), Inches(0.3),
            'LLM → Mirror KG 写入链路', font_size=Pt(16), color=DARK_BLUE, bold=True)

# Flow
write_boxes = [
    ('LLM 输出', Inches(0.7), DARK_BLUE),
    ('llm_extraction_runs\n+ items (raw_response)', Inches(3.2), MID_BLUE),
    ('llm_to_mirror\n(确定性转换)', Inches(6.8), ORANGE),
    ('Mirror KG\n(8 张表)', Inches(9.8), GREEN),
]
for name, wx, wc in write_boxes:
    flow_box(sl, wx, Inches(4.75), Inches(2.2), Inches(0.8), name, color=wc, font_size=Pt(11))
for wx in [Inches(2.95), Inches(6.55), Inches(9.55)]:
    flow_arrow(sl, wx, Inches(5.0), Inches(0.2), color=MID_GRAY)

# Mirror tables list
mirror_tables = 'mirror_region_connections  ·  mirror_region_functions  ·  mirror_region_circuits  ·  mirror_circuit_steps  ·  mirror_circuit_functions  ·  mirror_projection_functions  ·  mirror_kg_triples  ·  mirror_evidence_records'
add_textbox(sl, Inches(0.7), Inches(5.8), Inches(12), Inches(0.55),
            mirror_tables, font_size=Pt(10), color=MID_GRAY)
add_textbox(sl, Inches(0.7), Inches(6.2), Inches(12), Inches(0.3),
            'status: mirror_candidate → llm_suggested → rule_checked → human_review_pending',
            font_size=Pt(11), color=MID_GRAY)

# Bottom red banner
add_rect(sl, Inches(0.7), Inches(6.7), Inches(11.5), Inches(0.45), fill_color=RGBColor(0xFE, 0xF0, 0xF0))
add_textbox(sl, Inches(1.0), Inches(6.72), Inches(10.9), Inches(0.4),
            '⚠ LLM 输出绝不直接写入 final_*，不自动审核，不自动晋升 —— 所有 LLM 产出必须经过 Mirror KG 中转与人工审核',
            font_size=Pt(12), color=RED_ACCENT, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Mirror KG Governance
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, 'Mirror KG 治理机制')
slide_number(sl, 8)

# Why Mirror KG
add_textbox(sl, Inches(0.7), Inches(1.5), Inches(5), Inches(0.3),
            'Mirror KG = 预正式知识中转层', font_size=Pt(16), color=DARK_BLUE, bold=True)
add_textbox(sl, Inches(0.7), Inches(1.85), Inches(11.5), Inches(0.35),
            '解决 LLM 多 run / 多 pack 重叠提取 · 重跑版本差异 · 审核员重复数据困境',
            font_size=Pt(12), color=MID_GRAY)

# Left: Dedup Merge
process_card(sl, Inches(0.7), Inches(2.4), Inches(5.7), Inches(2.3),
             '写入时去重合并',
             ['每种实体定义 Canonical Key，写入时自动匹配',
              '',
              '▸ 连接: (source, target, type, directionality)，无向时排序',
              '▸ 回路: (circuit_name, source_atlas, granularity)',
              '▸ 回路步骤: (circuit_id, step_order)',
              '',
              '合并策略：高置信度胜出 + 双溯源保留',
              '不合并：已审核 / 已晋升 / 跨 atlas / 跨粒度'],
             color=MID_BLUE)

# Right: Dual Model Review
process_card(sl, Inches(6.8), Inches(2.4), Inches(5.7), Inches(2.3),
             '双模型盲审 (DeepSeek + Kimi)',
             ['同一数据项，两个模型独立审核，互相不可见对方结果',
              '',
              '  consensus → 双模型一致 — 加速审核通道',
              '  conflict  → 模型冲突 — 升级人工裁决',
              '',
              '冲突时提供差异化分析，标注两模型分歧点',
              '帮助审核员快速定位争议焦点'],
             color=ORANGE)

# Bottom: Cross Validation
process_card(sl, Inches(0.7), Inches(5.0), Inches(11.8), Inches(1.2),
             '回路-投射交叉验证（确定性算法，不调 LLM）',
             ['正向推导：回路 → 步骤 → 投射    ←→    反向聚合：投射 → 回路    →    交叉比对',
              '结果: bidirectionally_supported（通过）/ conflict（需人工），为人工审核提供结构化置信度信号'],
             color=GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Validation Center: Three Gates
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, '校验中心：三道闸门')
slide_number(sl, 9)

# Three gate columns
gate_data = [
    ('闸门 1\n规则校验', '12 规则\n确定性, 无 LLM', DARK_BLUE,
     ['Blocker → Tier1 自动修复 / Tier2 LLM 增强', 'Warning → 标记提醒审核员']),
    ('闸门 2\n大模型校验', 'DeepSeek+Kimi\n双模型盲审', MID_BLUE,
     ['consensus → 绿色通道加速', 'conflict → 升级 + 分歧标注']),
    ('闸门 3\n人工审核', '专家终审\n唯一终审权', ORANGE,
     ['approve → 晋升队列', 'reject → 退回 + 原因', 'request_changes → 修改']),
]

for i, (title, subtitle, color, details) in enumerate(gate_data):
    gx = Inches(0.7 + i * 4.2)
    # Gate box
    add_rect(sl, gx, Inches(1.6), Inches(3.8), Inches(3.7), fill_color=WHITE)
    # Title bar
    add_rect(sl, gx, Inches(1.6), Inches(3.8), Inches(1.2), fill_color=color)
    add_textbox(sl, gx + Inches(0.2), Inches(1.65), Inches(3.4), Inches(0.8),
                title, font_size=Pt(18), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(sl, gx + Inches(0.2), Inches(2.2), Inches(3.4), Inches(0.5),
                subtitle, font_size=Pt(11), color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    # Details
    dy = Inches(3.0)
    for d in details:
        add_textbox(sl, gx + Inches(0.3), dy, Inches(3.2), Inches(0.4),
                    '▸ ' + d, font_size=Pt(11), color=DARK_GRAY)
        dy += Inches(0.45)
    # Gate number circle
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, gx + Inches(1.5), Inches(3.7), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    cp = ctf.paragraphs[0]
    cp.text = str(i + 1)
    cp.font.size = Pt(20)
    cp.font.color.rgb = WHITE
    cp.font.bold = True
    cp.alignment = PP_ALIGN.CENTER

# Arrows between gates
for ax in [Inches(4.55), Inches(8.75)]:
    flow_arrow(sl, ax, Inches(3.9), Inches(0.25), color=MID_GRAY)

# Convergence to Final KG
add_rect(sl, Inches(3.5), Inches(5.8), Inches(6.2), Inches(0.55), fill_color=DARK_BLUE)
add_textbox(sl, Inches(3.5), Inches(5.82), Inches(6.2), Inches(0.5),
            '→ Final KG（三道全过，缺一不可）←',
            font_size=Pt(15), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
# Vertical connectors
for vx in [Inches(2.6), Inches(6.8), Inches(11.0)]:
    add_rect(sl, vx, Inches(5.35), Pt(3), Inches(0.5), fill_color=MID_GRAY)

# Principle
add_textbox(sl, Inches(0.7), Inches(6.7), Inches(11.5), Inches(0.3),
            '设计原则：前两道自动化（规则 + LLM 校验）+ 最后一道人工终审 = 效率与质量的最佳平衡',
            font_size=Pt(13), color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Promotion & Final KG Model
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, '晋升与 Final KG 模型')
slide_number(sl, 10)

# Promotion flow
add_textbox(sl, Inches(0.7), Inches(1.5), Inches(5), Inches(0.3),
            '晋升 (Promotion)', font_size=Pt(16), color=DARK_BLUE, bold=True)

promo_boxes = [
    ('Mirror KG\n(审核通过)', DARK_BLUE),
    ('预览→确认→执行\n(强确认, 不可逆)', ORANGE),
    ('Final KG\n(8 表映射)', GREEN),
]
px = Inches(0.7)
for name, pc in promo_boxes:
    flow_box(sl, px, Inches(1.95), Inches(3.5), Inches(0.7), name, color=pc, font_size=Pt(12))
    if name != 'Final KG\n(8 表映射)':
        flow_arrow(sl, px + Inches(3.6), Inches(2.15), Inches(0.2), color=pc)
    px += Inches(3.8)

# Final KG: Three layer model
add_textbox(sl, Inches(0.7), Inches(2.95), Inches(5), Inches(0.3),
            'Final KG 三元组模型', font_size=Pt(16), color=DARK_BLUE, bold=True)

# Three layers
layer_config = [
    ('实体层 (Nodes)', 'BrainRegion · Function · Circuit · Step · Projection', DARK_BLUE, Inches(2.5)),
    ('关系层 (Edges / Predicates)', '12 种标准谓词 (structurally_connects_to / has_function / has_step / involves_region / ...)', MID_BLUE, Inches(3.5)),
    ('统一查询层', 'final_kg_triples (subject, predicate, object) —— 确定性 Triple Consolidation', GREEN, Inches(4.5)),
]
for label, desc, lc, ly in layer_config:
    if ly == Inches(2.5):
        lh = Inches(0.55)
    elif ly == Inches(3.5):
        lh = Inches(0.55)
    else:
        lh = Inches(0.75)
    add_rect(sl, Inches(0.7), ly, Inches(7.5), lh, fill_color=lc)
    add_textbox(sl, Inches(1.0), ly + Pt(3), Inches(6.9), Inches(0.25),
                label, font_size=Pt(13), color=WHITE, bold=True)
    add_textbox(sl, Inches(1.0), ly + Inches(0.25), Inches(6.9), Inches(0.3),
                desc, font_size=Pt(9), color=RGBColor(0xEE, 0xEE, 0xFF))

# Granularity isolation
add_textbox(sl, Inches(9.0), Inches(2.95), Inches(4.0), Inches(0.3),
            '五层粒度隔离', font_size=Pt(16), color=DARK_BLUE, bold=True)

granularities = [
    ('macro_clinical', 'AAL3, Macro96', DARK_BLUE),
    ('meso_anatomical', 'HCP-MMP, Desikan', MID_BLUE),
    ('sub_connectivity', 'Brainnetome', RGBColor(0x2B, 0x6C, 0xB0)),
    ('fine_cyto', 'Julich-Brain', ORANGE),
    ('molecular_attr', 'Allen', GREEN),
]
gy = Inches(3.4)
for gname, gsrc, gc in granularities:
    add_textbox(sl, Inches(9.2), gy, Inches(1.8), Inches(0.25),
                gname, font_size=Pt(10), color=gc, bold=True)
    add_textbox(sl, Inches(11.0), gy, Inches(2.0), Inches(0.25),
                '← ' + gsrc, font_size=Pt(9), color=MID_GRAY)
    gy += Inches(0.28)
add_textbox(sl, Inches(9.0), Inches(4.9), Inches(4.0), Inches(0.4),
            '跨粒度: 显式 Mapping\n(exact_match / part_of / overlaps)\n禁止名称相似度自动合并',
            font_size=Pt(10), color=ORANGE, bold=True)

# Predicate table at bottom
add_textbox(sl, Inches(0.7), Inches(5.5), Inches(5), Inches(0.3),
            '12 种标准谓词', font_size=Pt(14), color=DARK_BLUE, bold=True)
add_table(sl, Inches(0.7), Inches(5.85),
          [Inches(3.5), Inches(3.5), Inches(2.7)],
          ['谓词 (Predicate)', '含义', '方向'],
          [
              ['structurally_connects_to', '结构连接', '脑区 → 脑区'],
              ['functionally_connects_to', '功能连接', '脑区 → 脑区'],
              ['projects_to', '投射', '脑区 → 脑区'],
              ['has_function', '区域功能', '脑区 → 功能'],
              ['has_step / involves_region', '回路步骤', '回路 → 步骤 → 脑区'],
              ['has_projection_function', '投射功能', '投射 → 功能'],
          ], font_size=Pt(10))


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Knowledge Consumption & Provenance
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, '知识消费与全链路溯源')
slide_number(sl, 11)

# Knowledge consumption
add_textbox(sl, Inches(0.7), Inches(1.5), Inches(5), Inches(0.3),
            '知识消费面', font_size=Pt(16), color=DARK_BLUE, bold=True)

consumers = [
    ('图谱探索', 'D3 力导向图，节点展开/收起\n按类型着色，置信度透明度', '🔍'),
    ('症状查询', '自然语言症状 → 标准化功能/回路\n→ 图谱结果 + 临床报告', '🏥'),
    ('数据中心', 'Raw → Candidate → Mirror → Final\n四面板浏览 + 字段补全', '📊'),
    ('知识导出', 'JSONL + CSV + Neo4j 兼容\n离线确定性导出', '📦'),
]
for i, (title, desc, icon) in enumerate(consumers):
    cx = Inches(0.7 + i * 3.2)
    process_card(sl, cx, Inches(1.95), Inches(3.0), Inches(1.6), icon + '  ' + title, [desc], color=MID_BLUE)

# Full provenance chain
add_textbox(sl, Inches(0.7), Inches(3.9), Inches(5), Inches(0.3),
            '全链路溯源（7 步回溯）', font_size=Pt(16), color=DARK_BLUE, bold=True)

provenance_chain = [
    ('Final KG 事实', DARK_BLUE),
    ('promotion_run', MID_BLUE),
    ('review_record', RGBColor(0x2B, 0x6C, 0xB0)),
    ('rule_validation_result', ORANGE),
    ('llm_extraction_item (raw_response)', RGBColor(0x80, 0x50, 0xB0)),
    ('llm_extraction_run (model + prompt)', RGBColor(0x6B, 0x46, 0xA0)),
    ('candidate_pool', RGBColor(0x2B, 0x6C, 0xB0)),
    ('import_batch', MID_BLUE),
    ('resource (原始脑图谱)', DARK_BLUE),
]

px_start = Inches(0.4)
pbox_w = Inches(1.3)
pbox_h = Inches(0.55)
py = Inches(4.4)

for i, (name, pc) in enumerate(provenance_chain):
    flow_box(sl, px_start, py, pbox_w, pbox_h, name, color=pc, font_size=Pt(9))
    if i < len(provenance_chain) - 1:
        flow_arrow(sl, px_start + pbox_w, py + Inches(0.13), Inches(0.13), color=MID_GRAY)
    px_start += pbox_w + Inches(0.16)

# Bottom: key message
add_rect(sl, Inches(0.7), Inches(5.3), Inches(11.5), Inches(0.5), fill_color=DARK_BLUE)
add_textbox(sl, Inches(1.0), Inches(5.33), Inches(10.9), Inches(0.45),
            '任何 Final KG 事实 → 7 步回溯 → 原始脑图谱资源：全链路不可篡改，provenance 是晋升的硬性前提',
            font_size=Pt(13), color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Evidence layer detail
add_textbox(sl, Inches(0.7), Inches(6.1), Inches(12), Inches(0.3),
            '证据层记录', font_size=Pt(14), color=DARK_BLUE, bold=True)
add_textbox(sl, Inches(0.7), Inches(6.45), Inches(12), Inches(0.4),
            '每个实体附 mirror_evidence_records → 记录 source_atlas · source_version · llm_run_id · review_record_id · promotion_run_id · raw_response · confidence',
            font_size=Pt(11), color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Innovation Summary
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
title_bar(sl, '核心创新点')
slide_number(sl, 12)

innovations = [
    ('1', '分层漏斗治理',
     '6 阶段流水线 (候选→校验→LLM→Mirror→Review→Final)\n每层明确写边界，LLM 是工具而非终审',
     DARK_BLUE),
    ('2', 'Mirror KG 中转层',
     '写入时去重合并 + Canonical Key 体系\n双溯源保留，已审核数据永不自动合并',
     MID_BLUE),
    ('3', '双模型盲审',
     'DeepSeek + Kimi 独立审核同一数据\n互不可见结果，冲突时人工裁决',
     RGBColor(0x80, 0x50, 0xB0)),
    ('4', '数据增强引擎',
     'Tier 1 确定性修复 (零 LLM 成本) + Tier 2 LLM 增强\nQ uality Score 0-100 自动评分',
     GREEN),
    ('5', '全链路溯源',
     '7 步回溯 Final KG → 原始脑图谱资源\n所有 provenance 不可变，是晋升硬性前提',
     ORANGE),
    ('6', '五层粒度隔离',
     'PostgreSQL schema 级物理隔离\n跨粒度显式 Mapping，禁止名称相似度合并',
     RED_ACCENT),
]

for i, (num, title, desc, color) in enumerate(innovations):
    row = i // 3
    col = i % 3
    ix = Inches(0.7 + col * 4.2)
    iy = Inches(1.6 + row * 2.8)

    # Card
    add_rect(sl, ix, iy, Inches(3.8), Inches(2.4), fill_color=WHITE)
    # Top accent bar
    add_rect(sl, ix, iy, Inches(3.8), Pt(5), fill_color=color)
    # Number circle
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, ix + Inches(0.2), iy + Inches(0.2), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    cp = ctf.paragraphs[0]
    cp.text = num
    cp.font.size = Pt(16)
    cp.font.color.rgb = WHITE
    cp.font.bold = True
    cp.alignment = PP_ALIGN.CENTER
    # Title
    add_textbox(sl, ix + Inches(0.9), iy + Inches(0.25), Inches(2.7), Inches(0.35),
                title, font_size=Pt(14), color=color, bold=True)
    # Description
    add_textbox(sl, ix + Inches(0.25), iy + Inches(0.85), Inches(3.3), Inches(1.3),
                desc, font_size=Pt(10), color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Status & Next Steps
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, LIGHT_GRAY)
title_bar(sl, '当前状态与下一步规划')
slide_number(sl, 13)

# Completed items
add_textbox(sl, Inches(0.7), Inches(1.5), Inches(5), Inches(0.3),
            '✅ 已实现', font_size=Pt(18), color=GREEN, bold=True)

completed = [
    '端到端流水线 (Resource → Final KG) 闭环',
    '7 种 LLM 提取能力 (连接/功能/回路/步骤/投射功能/回路功能/三元组)',
    '12 条确定性规则校验 + 数据增强引擎 (Tier 1 + Tier 2)',
    '双模型盲审 (DeepSeek + Kimi) + 交叉验证',
    'Mirror → Final 晋升 + Triple Consolidation',
    '写入时去重合并 (Canonical Key), 6 类实体全覆盖',
    '14 页前端工作台 + Graph Explorer + 症状查询 + 3D 脑区',
    '1,173 测试函数 · 59 数据库迁移 · 全链路 audit trail',
]
cy = Inches(2.0)
for item in completed:
    add_textbox(sl, Inches(1.0), cy, Inches(5.5), Inches(0.32),
                '▸ ' + item, font_size=Pt(12), color=DARK_GRAY)
    cy += Inches(0.35)

# Next steps
add_textbox(sl, Inches(7.0), Inches(1.5), Inches(5), Inches(0.3),
            '🚀 规划中', font_size=Pt(18), color=ORANGE, bold=True)

next_steps = [
    '接入更多粒度 (meso_anatomical, sub_connectivity, fine_cyto, molecular_attr)',
    '图数据库同步 (Neo4j 可选路径)',
    '跨粒度映射关系自动发现',
    'Graph Explorer 交互增强 (ReactFlow 迁移)',
    'DashBoard 关键指标看板',
    '批量字段补全进度可视化',
    '知识图谱版本管理与差异对比',
]
ny = Inches(2.0)
for item in next_steps:
    add_textbox(sl, Inches(7.3), ny, Inches(5.5), Inches(0.32),
                '▸ ' + item, font_size=Pt(12), color=DARK_GRAY)
    ny += Inches(0.35)

# Divider
add_rect(sl, Inches(6.6), Inches(1.6), Pt(2), Inches(5.2), fill_color=RGBColor(0xE2, 0xE8, 0xF0))

# Bottom metrics bar
metrics = [
    ('42', 'API 路由'),
    ('88', '服务模块'),
    ('14', '前端页面'),
    ('59', '数据库迁移'),
    ('1,173', '测试函数'),
    ('5', '粒度 Schema'),
]
mx = Inches(1.0)
for val, label in metrics:
    add_rect(sl, mx, Inches(6.3), Inches(1.75), Inches(0.85), fill_color=WHITE)
    add_textbox(sl, mx, Inches(6.35), Inches(1.75), Inches(0.4),
                val, font_size=Pt(22), color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(sl, mx, Inches(6.72), Inches(1.75), Inches(0.3),
                label, font_size=Pt(10), color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    mx += Inches(1.95)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You / Q&A
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, DARK_BLUE)

add_textbox(sl, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0),
            '感谢关注', font_size=Pt(48), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.6),
            '欢迎提问与交流', font_size=Pt(24), color=RGBColor(0xA0, 0xC4, 0xE8), alignment=PP_ALIGN.CENTER)

accent_line(sl, Inches(5), Inches(4.3), Inches(3.3))

add_textbox(sl, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.5),
            'NeuroGraphIQ KG V3 — 多粒度脑区知识图谱',
            font_size=Pt(16), color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.4),
            '2026 年 8 月', font_size=Pt(12), color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_rect(sl, Inches(1.5), Inches(7.0), Inches(10.3), Pt(1), fill_color=RGBColor(0x3A, 0x56, 0x7D))


# ── Save ───────────────────────────────────────────────────────
output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'docs', 'presentations')
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, 'NeuroGraphIQ_KG_V3_构建过程.pptx')
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')
